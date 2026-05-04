import os
import asyncio
import subprocess
import edge_tts
import nest_asyncio
import fitz  # PyMuPDF
from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips, CompositeAudioClip
import moviepy.audio.fx.all as afx
from pydub import AudioSegment
from pptx import Presentation

nest_asyncio.apply()

class PPTConverter:
    def __init__(self, ppt_path, bgm_path, silent_sec, libreoffice_path="soffice"):
        self.ppt_path = ppt_path
        self.bgm_path = bgm_path
        self.silent_sec = silent_sec
        self.libreoffice_path = libreoffice_path
        
        self.base_dir = os.path.dirname(ppt_path)
        self.base_name = os.path.splitext(os.path.basename(ppt_path))[0]
        self.pdf_path = os.path.join(self.base_dir, f"{self.base_name}.pdf")
        self.video_path = os.path.join(self.base_dir, f"{self.base_name}_output.mp4")
        self.audio_folder = os.path.join(self.base_dir, f"{self.base_name}_audio")
        self.image_folder = os.path.join(self.base_dir, f"{self.base_name}_images")

    def run_pipeline(self, progress_callback=None, video_logger=None):
        # 1. PPT 轉 PDF
        if progress_callback: progress_callback("正在將 PPT 轉換為 PDF...")
        self._convert_ppt_to_pdf()

        # 2. 提取備忘稿
        if progress_callback: progress_callback("正在提取備忘稿...")
        notes = self._extract_notes()

        # 3. 備忘稿轉語音
        if progress_callback: progress_callback("正在將備忘稿轉換為語音...")
        audio_paths = self._notes_to_audio(notes)

        # 4. PDF 轉圖片
        if progress_callback: progress_callback("正在將 PDF 轉換為圖片...")
        image_paths = self._pdf_to_images()

        # 5. 合成影片
        if progress_callback: progress_callback("正在合成最終影片，這可能需要一點時間...")
        # 這裡修改：把 video_logger 傳進去
        self._create_synced_video(image_paths, audio_paths, logger=video_logger)
        
        if progress_callback: progress_callback("轉換完成！")
        return self.video_path

    def _convert_ppt_to_pdf(self):
        # 呼叫本地端的 LibreOffice 進行轉檔
        subprocess.run([
            self.libreoffice_path, '--headless', '--convert-to', 'pdf', 
            '--outdir', self.base_dir, self.ppt_path
        ], check=True)

    def _extract_notes(self):
        prs = Presentation(self.ppt_path)
        all_notes = []
        for slide in prs.slides:
            text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                text = slide.notes_slide.notes_text_frame.text.strip()
            all_notes.append(text)
        return all_notes

    async def _generate_speech(self, text, output_file_path):
        if not text.strip():
            silent_ms = int(self.silent_sec * 1000)
            silence = AudioSegment.silent(duration=silent_ms)
            silence.export(output_file_path, format="wav")
        else:
            communicate = edge_tts.Communicate(text, "zh-TW-HsiaoChenNeural")
            await communicate.save(output_file_path)

    def _notes_to_audio(self, notes):
        os.makedirs(self.audio_folder, exist_ok=True)
        audio_paths = []
        for i, text in enumerate(notes):
            output_file_path = os.path.join(self.audio_folder, f"audio_segment_{i+1:03d}.wav")
            asyncio.run(self._generate_speech(text, output_file_path))
            audio_paths.append(output_file_path)
        return audio_paths

    def _pdf_to_images(self):
        pdf_document = fitz.open(self.pdf_path)
        os.makedirs(self.image_folder, exist_ok=True)
        image_paths = []
        for page_number in range(len(pdf_document)):
            page = pdf_document.load_page(page_number)
            pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
            image_path = os.path.join(self.image_folder, f"temp_page_{page_number}.png")
            pix.save(image_path)
            image_paths.append(image_path)
        return image_paths

    def _create_synced_video(self, image_paths, audio_paths, logger=None):
        clips = []
        for img_path, au_path in zip(image_paths, audio_paths):
            audio_clip = AudioFileClip(au_path)
            img_clip = ImageClip(img_path).set_duration(audio_clip.duration).resize(height=1080)
            img_clip = img_clip.set_position('center').on_color(size=(1920, 1080), color=(0,0,0), col_opacity=1)
            img_clip = img_clip.set_audio(audio_clip)
            clips.append(img_clip)

        final_video = concatenate_videoclips(clips, method="chain")

        if self.bgm_path and os.path.exists(self.bgm_path):
            bgm_clip = AudioFileClip(self.bgm_path)
            bgm_clip = afx.audio_loop(bgm_clip, duration=final_video.duration).volumex(0.5)
            mixed_audio = CompositeAudioClip([final_video.audio, bgm_clip])
            final_video = final_video.set_audio(mixed_audio)

        final_video.write_videofile(
            self.video_path, fps=24, codec="libx264", audio_codec="aac",
            threads=4, preset='ultrafast', ffmpeg_params=["-pix_fmt", "yuv420p"],
            logger=logger  # <--- 這裡修改：把 logger 參數交給 moviepy
        )
